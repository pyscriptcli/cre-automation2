import os
import io
import subprocess
import tempfile
import streamlit as st
from pptx import Presentation
from PIL import Image

# --- FLAT & LUXURIOUS LUXURY DESIGN SYSTEM (CSS INJECTION) ---
LUXURY_CRE_SYSTEM = """
<style>
    /* 1. Global Page Core Reset to White & Navy */
    .stApp {
        background-color: #FFFFFF !important;
        color: #002B49 !important;
        font-family: 'Inter', -apple-system, sans-serif !important;
    }
    
    /* Remove all default header and padding layers */
    div[data-testid="stHeader"] {
        background-color: #FFFFFF !important;
    }
    .block-container {
        padding-top: 2rem !important;
        padding-bottom: 2rem !important;
        max-width: 1200px !important; /* Expanded width for side-by-side single page execution */
    }
    
    /* 2. Flat Luxury Input Fields & Dropdowns (Larger Fonts, Sharp 0px Corners) */
    div[data-baseweb="input"], div[data-baseweb="base-input"], div[role="textbox"], div[data-baseweb="select"] {
        background-color: #FFFFFF !important;
        border: 1px solid #002B49 !important; /* Solid Navy Border Line */
        border-radius: 0px !important; /* Flat luxury edge */
        color: #002B49 !important;
        transition: border-color 0.15s ease-in-out !important;
    }
    div[data-baseweb="input"]:focus-within, div[data-baseweb="select"]:focus-within {
        border-color: #C5A059 !important; /* Corporate Gold Accent */
        box-shadow: none !important;
    }
    
    /* Scaled Input text parameters */
    input[type="text"], .stTextInput input, div[data-baseweb="select"] div {
        color: #002B49 !important;
        background-color: #FFFFFF !important;
        font-size: 16px !important; /* Larger font size override */
        font-weight: 500 !important;
    }
    input[type="text"], .stTextInput input {
        padding: 10px 14px !important;
    }
    
    /* Luxury Corporate Checkbox Overrides */
    div[data-testid="stCheckbox"] label p {
        color: #002B49 !important;
        font-size: 14px !important;
        font-weight: 600 !important;
    }
    
    /* 3. Minimalist Template File Uploader */
    section[data-testid="stFileUploader"] {
        background-color: #FFFFFF !important;
        border: 1px solid #002B49 !important; /* Matches text rows exactly */
        border-radius: 0px !important;
        padding: 4px 12px !important;
    }
    section[data-testid="stFileUploader"] div, section[data-testid="stFileUploader"] span {
        color: #002B49 !important;
        font-size: 13px !important;
        font-weight: 500 !important;
    }
    div[data-testid="stFileUploader"] label {
        display: none !important; /* Remove native duplicate labels */
    }
    
    /* 4. Large Bold Typography Row Layouts */
    .row-metric-label {
        font-size: 15px !important; /* Scaled up text label */
        font-weight: 700 !important;
        text-transform: uppercase !important;
        letter-spacing: 0.08em !important;
        color: #002B49 !important;
        display: flex;
        align-items: center;
        padding-top: 12px;
    }
    .large-icon {
        font-size: 24px !important; /* Magnified flat icon framework */
        margin-right: 12px;
        display: inline-block;
    }
    
    /* 5. Clean Section Framing */
    .luxury-workspace-card {
        background-color: #FFFFFF;
        border-top: 4px solid #002B49; /* Massive Navy Block Line */
        padding-top: 20px;
        margin-bottom: 25px;
    }

    /* 6. Heavy Command Button Overrides */
    div.stButton > button {
        background-color: #002B49 !important;
        color: #FFFFFF !important;
        font-weight: 700 !important;
        font-size: 14px !important;
        text-transform: uppercase !important;
        letter-spacing: 0.08em !important;
        border: none !important;
        border-radius: 0px !important;
        border-bottom: 4px solid #C5A059 !important; /* Prominent Gold Base Accent */
        padding: 16px 32px !important;
        width: 100% !important;
        transition: background-color 0.15s ease;
    }
    div.stButton > button:hover {
        background-color: #0A3352 !important;
        border-bottom-color: #C5A059 !important;
        color: #FFFFFF !important;
    }
    
    /* Reset Button alternative alignment */
    div[data-testid="column"] div.stButton > button[key="reset_btn_key"] {
        background-color: transparent !important;
        color: #64748B !important;
        border: 1px solid #CBD5E1 !important;
        border-bottom: 1px solid #CBD5E1 !important;
        padding: 12px 24px !important;
    }
    div[data-testid="column"] div.stButton > button[key="reset_btn_key"]:hover {
        color: #002B49 !important;
        border-color: #002B49 !important;
    }
</style>
"""

# --- BACKEND SPATIAL PROCESSING UTILS ---
def smart_crop_to_fit(img_file, target_w_emu, target_h_emu):
    try:
        img = Image.open(img_file)
        img_w, img_h = img.size
        target_ratio = target_w_emu / target_h_emu
        img_ratio = img_w / img_h
        
        if img_ratio > target_ratio:
            new_w = int(img_h * target_ratio)
            left = (img_w - new_w) // 2
            right = left + new_w
            img = img.crop((left, 0, right, img_h))
        else:
            new_h = int(img_w / target_ratio)
            top = (img_h - new_h) // 2
            bottom = top + new_h
            img = img.crop((0, top, img_w, bottom))
            
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        return img_byte_arr
    except Exception:
        return img_file

# --- HEADLESS LINUX PDF GENERATION MATRIX ---
def convert_pptx_to_pdf(pptx_bytes):
    with tempfile.TemporaryDirectory() as temp_dir:
        input_pptx_path = os.path.join(temp_dir, "document.pptx")
        with open(input_pptx_path, "wb") as f:
            f.write(pptx_bytes)
        try:
            command = ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", temp_dir, input_pptx_path]
            subprocess.run(command, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            output_pdf_path = os.path.join(temp_dir, "document.pdf")
            if os.path.exists(output_pdf_path):
                with open(output_pdf_path, "rb") as f:
                    return f.read()
        except Exception:
            return None
    return None

# --- APP CONFIGURATION ---
st.set_page_config(page_title="Asset Engine", page_icon="🏢", layout="wide")
st.markdown(LUXURY_CRE_SYSTEM, unsafe_allow_html=True)

# HIGH DENSITY STRUCTURAL ROW HELPERS
def dynamic_form_row(icon, label_text, default_val, state_key):
    r_col1, r_col2 = st.columns([9, 11])
    with r_col1:
        st.markdown(f'<div class="row-metric-label"><span class="large-icon">{icon}</span> {label_text}</div>', unsafe_allow_html=True)
    with r_col2:
        return st.text_input("", value=default_val, key=state_key, label_visibility="collapsed")

def dynamic_uploader_row(icon, label_text, allowed_types, state_key):
    r_col1, r_col2 = st.columns([9, 11])
    with r_col1:
        st.markdown(f'<div class="row-metric-label"><span class="large-icon">{icon}</span> {label_text}</div>', unsafe_allow_html=True)
    with r_col2:
        return st.file_uploader(label_text, type=allowed_types, key=state_key, label_visibility="collapsed")

def dynamic_selector_row(icon, label_text, options, state_key):
    r_col1, r_col2 = st.columns([9, 11])
    with r_col1:
        st.markdown(f'<div class="row-metric-label"><span class="large-icon">{icon}</span> {label_text}</div>', unsafe_allow_html=True)
    with r_col2:
        return st.selectbox(label_text, options, key=state_key, label_visibility="collapsed")

# --- TWO COLUMN MASTER WORKSPACE (ONE PAGE MATRIX) ---
col_left, col_right = st.columns([1, 1], gap="large")

with col_left:
    st.markdown('<div class="luxury-workspace-card">', unsafe_allow_html=True)
    prop_location = dynamic_form_row("📍", "Property Location", "Tagaytay, Cavite", "cre_loc")
    prop_size     = dynamic_form_row("📐", "Property Size (SQM)", "386", "cre_size")
    prop_type     = dynamic_form_row("🏢", "Property Type", "Commercial Space", "cre_type")
    prop_address  = dynamic_form_row("🗺️", "Full Address", "Mendez Crossing East, Tagaytay City, Cavite", "cre_addr")
    lease_rates   = dynamic_form_row("💰", "Lease Rates", "200,000 per month", "cre_rates")
    sec_deposit   = dynamic_form_row("🛡️", "Security Deposit", "3 months", "cre_sec")
    adv_rent      = dynamic_form_row("💵", "Advance Rent", "3 months", "cre_adv")
    escalation    = dynamic_form_row("📈", "Rental Escalation", "5%", "cre_esc")
    lease_term    = dynamic_form_row("📅", "Lease Term", "5 years", "cre_term")
    handover      = dynamic_form_row("🏗️", "Handover Condition", "As is where is", "cre_hand")
    
    # --- NEW: CALL TO ACTION SECTION (DYNAMIC MAX-2 CHECKBOX MATRIX) ---
    st.markdown("<div style='margin-top: 15px;'></div>", unsafe_allow_html=True)
    cta_col1, cta_col2 = st.columns([9, 11])
    with cta_col1:
        st.markdown('<div class="row-metric-label"><span class="large-icon">📞</span> CALL TO ACTION</div>', unsafe_allow_html=True)
    with cta_col2:
        contacts_database = {
            "Sondi Tuazon": {"phone": "0917 843 6128", "email": "sondi.tuazon@primephilippines.com"},
            "Meliza Zapata": {"phone": "0917 555 1234", "email": "meliza.zapata@primephilippines.com"},
            "Dykstra Pineda": {"phone": "0917 555 5678", "email": "dykstra.pineda@primephilippines.com"},
            "Cedtrix Rena": {"phone": "0917 555 9012", "email": "cedtriz.rena@primephilippines.com"},
            "Carlo Medina": {"phone": "0920 986 2764", "email": "carlo.medina@primephilippines.com"},
            "Dave Policarpio": {"phone": "0917 555 3456", "email": "dave.policarpio@primephilippines.com"},
            "Irish Rima": {"phone": "0917 555 7890", "email": "irish.rima@primephilippines.com"}
        }
        
        # Calculate currently selected count ahead of layout pass to enforce freezing
        selected_names = [name for name in contacts_database if st.session_state.get(f"state_cb_{name}", False)]
        current_checked_count = len(selected_names)
        
        # Split contacts array into a beautiful internal two-column layout
        cb_sub_col1, cb_sub_col2 = st.columns(2)
        contact_keys = list(contacts_database.keys())
        
        for idx, name in enumerate(contact_keys):
            is_active = st.session_state.get(f"state_cb_{name}", False)
            # Freeze rule: Disable if 2 elements are checked, unless it is the element currently active
            should_freeze = (current_checked_count >= 2 and not is_active)
            
            target_sub_column = cb_sub_col1 if idx % 2 == 0 else cb_sub_col2
            with target_sub_column:
                st.checkbox(name, key=f"state_cb_{name}", disabled=should_freeze)
                
    st.markdown('</div>', unsafe_allow_html=True)

with col_right:
    st.markdown('<div class="luxury-workspace-card">', unsafe_allow_html=True)
    t_col1, t_col2 = st.columns([9, 11])
    with t_col1:
        st.markdown('<div class="row-metric-label"><span class="large-icon">📂</span> TEMPLATE</div>', unsafe_allow_html=True)
    with t_col2:
        u_template = st.file_uploader("Template File Input", type=["pptx"], key="web_tpl")
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="luxury-workspace-card">', unsafe_allow_html=True)
    img_types = ["png", "jpg", "jpeg"]
    u_photo1  = dynamic_uploader_row("📸", "Property Photo 1", img_types, "web_p1")
    u_map     = dynamic_uploader_row("🗺️", "Location Map", img_types, "web_mp")
    u_lotplan = dynamic_uploader_row("📐", "Lot Plan", img_types, "web_lp")
    u_photo2  = dynamic_uploader_row("📸", "Property Photo 2", img_types, "web_p2")
    u_photo3  = dynamic_uploader_row("📸", "Property Photo 3", img_types, "web_p3")
    st.markdown('</div>', unsafe_allow_html=True)

    st.markdown('<div class="luxury-workspace-card">', unsafe_allow_html=True)
    export_format = dynamic_selector_row("📄", "Export Format", ["PPTX (PowerPoint)", "PDF (Adobe Acrobat)"], "web_format")
    st.markdown('</div>', unsafe_allow_html=True)

# Compile text data tokens
data_inputs = {
    "{{PROPERTY_LOCATION}}": prop_location, "{{PROPERTY_SIZE}}": prop_size,
    "{{PROPERTY_TYPE}}": prop_type, "{{PROPERTY_ADDRESS}}": prop_address,
    "{{LEASE_RATES}}": lease_rates, "{{SECURITY_DEPOSIT}}": sec_deposit,
    "{{ADVANCE_RENT}}": adv_rent, "{{ESCALATION}}": escalation,
    "{{LEASE TERM}}": lease_term, "{{HANDOVER CONDITION}}": handover
}

# --- FIXED: CONTACT TOKEN DATA PIPELINE MAPPING ---
# Add Contact Tokens dynamically based on active selected layout slots
for contact_slot in range(2):
    slot_num = contact_slot + 1
    
    # Construct tokens using explicit string concatenation to bypass f-string parser limits
    name_token  = "{{CONTACT_NAME_"  + str(slot_num) + "}}"
    phone_token = "{{CONTACT_PHONE_" + str(slot_num) + "}}"
    email_token = "{{CONTACT_EMAIL_" + str(slot_num) + "}}"
    
    if contact_slot < len(selected_names):
        target_name = selected_names[contact_slot]
        data_inputs[name_token]  = target_name
        data_inputs[phone_token] = contacts_database[target_name]["phone"]
        data_inputs[email_token] = contacts_database[target_name]["email"]
    else:
        # Graceful fallbacks for empty selections to clear remaining slide tags
        data_inputs[name_token]  = ""
        data_inputs[phone_token] = ""
        data_inputs[email_token] = ""
        
# --- CONTROL DESK ACTION LAYER ---
st.markdown("<div style='margin-top: 10px; border-top: 1px solid #002B49; padding-top: 20px;'></div>", unsafe_allow_html=True)
action_col1, action_col2 = st.columns([1, 2])

with action_col1:
    if st.button("↺ Clear", key="reset_btn_key", use_container_width=True):
        st.cache_data.clear()
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()

with action_col2:
    if u_template is None:
        st.markdown("<div style='padding-top:16px; font-size:12px; font-weight:700; color:#002B49; text-align:right; letter-spacing:0.04em;'>⚠️ UPLOAD A MASTER PPTX BLUEPRINT TO MOUNT GENERATION FUNCTIONS.</div>", unsafe_allow_html=True)
    else:
        if st.button("⚙️ GENERATE DECK", key="generate_btn_key", use_container_width=True):
            with st.spinner("Processing template slide assets and applying smart crops..."):
                try:
                    prs = Presentation(u_template)
                    
                    for slide in prs.slides:
                        shapes_to_delete = []
                        images_to_add = []

                        for shape in slide.shapes:
                            if shape.has_text_frame and not any(img_token in shape.text_frame.text for img_token in image_inputs):
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        for token, value in data_inputs.items():
                                            if token in run.text:
                                                run.text = run.text.replace(token, value)

                            if shape.has_table:
                                for row in shape.table.rows:
                                    for cell in row.cells:
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                for token, value in data_inputs.items():
                                                    if token in run.text:
                                                        run.text = run.text.replace(token, value)

                            if shape.has_text_frame:
                                text_content = shape.text_frame.text
                                for img_token, img_file in image_inputs.items():
                                    if img_token in text_content and img_file is not None:
                                        images_to_add.append((img_file, shape.left, shape.top, shape.width, shape.height))
                                        shapes_to_delete.append(shape)

                        for img_file, left, top, width, height in images_to_add:
                            processed_img = smart_crop_to_fit(img_file, width, height)
                            slide.shapes.add_picture(processed_img, left, top, width=width, height=height)

                        for old_shape in shapes_to_delete:
                            sp = old_shape._element
                            sp.getparent().remove(sp)

                    pptx_stream = io.BytesIO()
                    prs.save(pptx_stream)
                    raw_pptx_bytes = pptx_stream.getvalue()
                    
                    st.markdown("""<div style="border-left: 4px solid #C5A059; background-color: #FFFFFF; padding: 16px; border-top: 1px solid #002B49; border-right: 1px solid #002B49; border-bottom: 1px solid #002B49; margin-top: 20px; text-align: center; color: #002B49; font-weight: 700; font-size: 13px; letter-spacing: 0.05em;">🎉 PRESENTATION COMPILED SUCCESSFULLY! READY FOR PRODUCTION DOWNLOAD.</div>""", unsafe_allow_html=True)
                    st.markdown("<div style='margin-top:12px;'></div>", unsafe_allow_html=True)
                    
                    safe_filename = f"PIS_{prop_location.replace(' ', '_')}"
                    
                    if "PDF" in export_format:
                        with st.spinner("Converting OpenXML geometries to vector PDF binary streams..."):
                            pdf_bytes = convert_pptx_to_pdf(raw_pptx_bytes)
                        if pdf_bytes:
                            st.download_button(
                                label="📥 DOWNLOAD BRANDED PDF DECK",
                                data=pdf_bytes,
                                file_name=f"{safe_filename}.pdf",
                                mime="application/pdf",
                                use_container_width=True
                            )
                    else:
                        st.download_button(
                            label="📥 DOWNLOAD BRANDED PPTX DECK",
                            data=raw_pptx_bytes,
                            file_name=f"{safe_filename}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )
                    
                except Exception as e:
                    st.error(f"Compilation core failure log description: {str(e)}")
