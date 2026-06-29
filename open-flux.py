import os
import io
import subprocess
import tempfile
import re
import json
import math
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from PIL import Image, ImageDraw
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import base64
import traceback
import time

# --- MAP SPECIFIC DEPENDENCIES ---
import folium
from folium.plugins import Draw
from streamlit_folium import st_folium
import requests

# --- CONTACTS DATABASE FOR CTA PRESETS ---
contacts_database = {
    "Sondi Tuazon": {"phone": "0917 843 6128", "email": "sondi.tuazon@primephilippines.com"},
    "Meliza Zapata": {"phone": "0996 880 5399", "email": "meliza.zapata@primephilippines.com"},
    "Dykstra Pineda": {"phone": "0920 986 2748", "email": "dykstra.pineda@primephilippines.com"},
    "Cedtrix Rena": {"phone": "0977 653 1494", "email": "cedtriz.rena@primephilippines.com"},
    "Carlo Medina": {"phone": "0920 986 2763", "email": "carlo.medina@primephilippines.com"},
    "Dave Policarpio": {"phone": "0908 865 8945", "email": "dave.policarpio@primephilippines.com"},
    "Irish Rima": {"phone": "0918 622 5346", "email": "irish.rima@primephilippines.com"}
}

# --- PROGRAMMATIC LIGHT MODE LOCK ---
_config_dir = ".streamlit"
_config_file = os.path.join(_config_dir, "config.toml")
if not os.path.exists(_config_file):
    os.makedirs(_config_dir, exist_ok=True)
    with open(_config_file, "w", encoding="utf-8") as f:
        f.write("[theme]\nbase=\"light\"\n")

# --- MINIMAL UI CSS ---
MINIMAL_CRE_SYSTEM = """
<style>
    #MainMenu {visibility: hidden;}
    header {visibility: hidden;}
    .stApp { margin-top: -50px; }
    .stDeployButton {display: none;}
    .stStatusWidget {display: none;}
    
    .stApp { background-color: #FFFFFF !important; color: #1A1A1A !important; font-family: 'Segoe UI', Arial, sans-serif !important; }
    div[data-testid="stHeader"] { background-color: #FFFFFF !important; display: none !important; }
    .block-container { padding-top: 0.5rem !important; padding-bottom: 0.5rem !important; max-width: 1300px !important; }
    
    div[data-baseweb="input"], div[data-baseweb="base-input"], div[role="textbox"], div[data-baseweb="select"], textarea {
        background-color: #FFFFFF !important; border: 1px solid #CCCCCC !important; border-radius: 4px !important;
        color: #1A1A1A !important;
    }
    div[data-baseweb="input"]:focus-within, div[data-baseweb="select"]:focus-within, textarea:focus { border-color: #003366 !important; box-shadow: none !important; }
    input[type="text"], .stTextInput input, div[data-baseweb="select"] div, textarea { color: #1A1A1A !important; font-size: 14px !important; }
    
    div[data-baseweb="select"] { min-height: 32px !important; }
    div[data-baseweb="select"] > div { min-height: 32px !important; padding: 0 8px !important; }
    div[data-baseweb="select"] select { font-size: 13px !important; padding: 2px 8px !important; }
    
    section[data-testid="stFileUploader"] { background-color: #F8F8F8 !important; border: 1px solid #CCCCCC !important; border-radius: 4px !important; padding: 4px 12px !important; }
    .workspace-card { background-color: #FFFFFF; border: 1px solid #E0E0E0; border-radius: 4px; padding: 16px; margin-bottom: 12px; }
    .editor-card { background-color: #F8F9FA; border: 2px solid #003366; border-radius: 6px; padding: 20px; }
    
    div.stButton > button { 
        background-color: #003366 !important; color: #FFFFFF !important; font-weight: 600 !important; font-size: 12px !important; 
        border: none !important; border-radius: 4px !important; padding: 6px 14px !important; width: 100% !important; min-height: 32px !important;
    }
    div.stButton > button:hover { background-color: #002244 !important; transform: translateY(-1px); box-shadow: 0 4px 8px rgba(0, 51, 102, 0.2); }
    div.stButton > button:disabled { background-color: #666666 !important; opacity: 0.6; cursor: not-allowed; }
    
    .field-label { font-size: 13px !important; font-weight: 600 !important; color: #1A1A1A !important; padding-top: 6px; }
    .section-header { font-size: 15px !important; font-weight: 700 !important; color: #1A1A1A !important; margin-bottom: 10px; }
    .saved-indicator { background-color: #E8F5E9; padding: 6px 12px; border-radius: 4px; font-size: 13px; color: #2E7D32; border-left: 3px solid #2E7D32; margin-top: 6px; }
    hr { margin: 12px 0 !important; border-color: #E0E0E0 !important; }
    
    div[data-testid="stForm"] { border: 1px solid #E0E0E0 !important; border-radius: 6px !important; padding: 1rem !important; background-color: #FFFFFF; }
    
    .placeholder-label {
        font-weight: 600 !important;
        font-size: 13px !important;
        color: #1A1A1A !important;
        margin-bottom: 4px;
    }
</style>
"""

# --- FILE MANAGEMENT FUNCTIONS ---
def get_storage_dir():
    storage_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "stored_templates")
    os.makedirs(storage_dir, exist_ok=True)
    return storage_dir

def get_temp_config_path(template_name):
    storage_dir = get_storage_dir()
    safe_name = re.sub(r'[^\w\-_. ]', '_', template_name or "unsaved_template")
    return os.path.join(storage_dir, f"{safe_name}_temp_form_data.json")

def get_github_templates():
    root_dir = os.path.dirname(os.path.abspath(__file__))
    templates = []
    if os.path.exists(root_dir):
        for file in os.listdir(root_dir):
            if file.startswith('template_') and (file.endswith('.pptx') or file.endswith('.docx')):
                filepath = os.path.join(root_dir, file)
                stat = os.stat(filepath)
                display_name = file.replace('template_', '').replace('.pptx', '').replace('.docx', '')
                templates.append({
                    'name': file, 'display_name': display_name, 'path': filepath,
                    'size': stat.st_size, 'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'type': 'PPTX' if file.endswith('.pptx') else 'DOCX', 'source': 'github'
                })
    return templates

def save_template_to_file(template_bytes, template_name):
    storage_dir = get_storage_dir()
    safe_name = re.sub(r'[^\w\-_. ]', '_', template_name)
    if not safe_name.endswith('.pptx') and not safe_name.endswith('.docx'):
        safe_name += '.docx'
    filepath = os.path.join(storage_dir, safe_name)
    with open(filepath, 'wb') as f:
        f.write(template_bytes)
    return filepath

def load_template_from_file(template_name):
    root_dir = os.path.dirname(os.path.abspath(__file__))
    root_filepath = os.path.join(root_dir, template_name)
    if os.path.exists(root_filepath):
        with open(root_filepath, 'rb') as f: return f.read()
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        with open(filepath, 'rb') as f: return f.read()
    return None

def get_saved_templates():
    storage_dir = get_storage_dir()
    templates = []
    if os.path.exists(storage_dir):
        for file in os.listdir(storage_dir):
            if file.endswith('.pptx') or file.endswith('.docx'):
                filepath = os.path.join(storage_dir, file)
                stat = os.stat(filepath)
                templates.append({
                    'name': file, 'display_name': file.replace('.pptx', '').replace('.docx', ''),
                    'path': filepath, 'size': stat.st_size, 'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'type': 'PPTX' if file.endswith('.pptx') else 'DOCX', 'source': 'stored'
                })
    templates.extend(get_github_templates())
    return templates

def delete_template_file(template_name):
    root_dir = os.path.dirname(os.path.abspath(__file__))
    if os.path.exists(os.path.join(root_dir, template_name)):
        st.warning("Cannot delete GitHub repository templates")
        return False
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        os.remove(filepath)
        config_name = template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        config_path = os.path.join(storage_dir, config_name)
        if os.path.exists(config_path): os.remove(config_path)
        temp_config = get_temp_config_path(template_name)
        if os.path.exists(temp_config): os.remove(temp_config)
        return True
    return False

def save_config_to_file(config_data, config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(config_data, f, indent=4)
    return filepath

def load_config_from_file(config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    if os.path.exists(filepath):
        with open(filepath, 'r', encoding='utf-8') as f: return json.load(f)
    return None

def auto_save_config():
    if st.session_state.saved_template_name and st.session_state.custom_mapping:
        config_name = st.session_state.saved_template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        save_config_to_file(st.session_state.custom_mapping, config_name)

# --- CTA PRESET FUNCTIONS ---
def detect_cta_sets():
    cta_sets = {}
    for token in st.session_state.tokens:
        clean_label = token.replace("{", "").replace("}", "").upper()
        match = re.match(r'CTA(\d+)_(NAME|CONTACT_NUMBER|EMAIL)', clean_label)
        if match:
            cta_num = int(match.group(1))
            field_type = match.group(2)
            if cta_num not in cta_sets:
                cta_sets[cta_num] = {'tokens': {}, 'fields': set()}
            cta_sets[cta_num]['tokens'][field_type] = token
            cta_sets[cta_num]['fields'].add(field_type)
    return cta_sets

def apply_cta_preset_autofill(cta_num, advisor_name):
    if advisor_name not in contacts_database: return False
    contact_info = contacts_database[advisor_name]
    cta_sets = detect_cta_sets()
    if cta_num not in cta_sets: return False
    tokens = cta_sets[cta_num]['tokens']
    
    current_values = {}
    for token in st.session_state.tokens:
        val_key = f"val_{token}"
        if val_key in st.session_state: current_values[token] = st.session_state[val_key]
    
    if 'NAME' in tokens:
        st.session_state[f"val_{tokens['NAME']}"] = advisor_name
        st.session_state.temp_form_data[tokens['NAME']] = advisor_name
    if 'CONTACT_NUMBER' in tokens:
        st.session_state[f"val_{tokens['CONTACT_NUMBER']}"] = contact_info["phone"]
        st.session_state.temp_form_data[tokens['CONTACT_NUMBER']] = contact_info["phone"]
    if 'EMAIL' in tokens:
        st.session_state[f"val_{tokens['EMAIL']}"] = contact_info["email"]
        st.session_state.temp_form_data[tokens['EMAIL']] = contact_info["email"]
    
    cta_tokens = set(tokens.values())
    for token, value in current_values.items():
        val_key = f"val_{token}"
        if token not in cta_tokens:
            if val_key not in st.session_state or st.session_state[val_key] != value:
                st.session_state[val_key] = value
                st.session_state.temp_form_data[token] = value
        
    if st.session_state.saved_template_name:
        temp_path = get_temp_config_path(st.session_state.saved_template_name)
        try:
            with open(temp_path, 'w', encoding='utf-8') as f: json.dump(st.session_state.temp_form_data, f, indent=4)
            return True
        except Exception: return False
    return True

# --- BASEMAP CONFIGURATION WITH IMPROVED RELIABILITY ---
BASEMAP_CONFIG = {
    "Satellite (Labels)": {"urls": ["https://mt1.google.com/vt/lyrs=y&x={x}&y={y}&z={z}"], "attribution": "Google"},
    "Satellite (Streets)": {"urls": ["https://mt1.google.com/vt/lyrs=y&x={x}&y={y}&z={z}&apistyle=s.t%3A2%7Cp.v%3Aoff"], "attribution": "Google"},
    "Satellite (Clean)": {"urls": ["https://mt1.google.com/vt/lyrs=s&x={x}&y={y}&z={z}"], "attribution": "Google"},
    "Street Map": {"urls": ["https://mt1.google.com/vt/lyrs=m&x={x}&y={y}&z={z}"], "attribution": "Google"},
    "OSM Carto Light": {"urls": ["https://a.basemaps.cartocdn.com/light_all/{z}/{x}/{y}.png"], "attribution": "CartoDB"},
    "Open Street Map": {"urls": ["https://a.tile.openstreetmap.org/{z}/{x}/{y}.png"], "attribution": "OpenStreetMap"}
}

def get_tile_urls(style_name):
    config = BASEMAP_CONFIG.get(style_name)
    return config["urls"] if config else BASEMAP_CONFIG["Street Map"]["urls"]

def get_attribution(style_name):
    config = BASEMAP_CONFIG.get(style_name)
    return config["attribution"] if config else ""

def fetch_tile_with_retry(url_template, zoom, x, y, headers, max_retries=3):
    for attempt in range(max_retries):
        url = url_template.format(z=zoom, x=x, y=y)
        try:
            resp = requests.get(url, headers=headers, timeout=10)
            if resp.status_code == 200: return resp.content
        except Exception: continue
    return None

# --- MAP BOUNDING BOX GENERATOR ---
def generate_static_map_bounds(n, s, e, w, pin_lat, pin_lon, style="Satellite (Streets)", pin_color="#003366", pin_size=18):
    def calculate_1km_bounds(lat, lon):
        lat_deg_per_km = 1.0 / 111.32
        lon_deg_per_km = 1.0 / (111.32 * math.cos(math.radians(lat)))
        return lat + lat_deg_per_km * 0.5, lat - lat_deg_per_km * 0.5, lon + lon_deg_per_km * 0.5, lon - lon_deg_per_km * 0.5
    
    bounds_valid = all(x is not None for x in [n, s, e, w])
    if bounds_valid and (abs(n - s) < 0.0001 or abs(e - w) < 0.0001): bounds_valid = False
    if not bounds_valid: n, s, e, w = calculate_1km_bounds(pin_lat, pin_lon)
    
    lon_span = e - w
    lat_span = n - s
    zoom = max(13, min(20, int(math.log2((360.0 / (lon_span if lon_span > 0 else 0.001)) * 8))))
    if lon_span < 0.01 and lat_span < 0.01: zoom = min(20, zoom + 2)
    
    def deg2num(lat_deg, lon_deg, z):
        lat_rad = math.radians(lat_deg)
        n_tiles = 2.0 ** z
        return int((lon_deg + 180.0) / 360.0 * n_tiles), int((1.0 - math.log(math.tan(lat_rad) + (1 / math.cos(lat_rad))) / math.pi) / 2.0 * n_tiles)
    
    x_min, y_min = deg2num(n, w, zoom)
    x_max, y_max = deg2num(s, e, zoom)
    if x_max == x_min: x_max += 1
    if y_max == y_min: y_max += 1
    
    tile_size, scale_factor = 256, 2
    stitched = Image.new('RGB', ((x_max - x_min + 1) * tile_size * scale_factor, (y_max - y_min + 1) * tile_size * scale_factor))
    headers = {"User-Agent": "Mozilla/5.0"}
    tile_urls = get_tile_urls(style)
    
    for x in range(x_min, x_max + 1):
        for y in range(y_min, y_max + 1):
            for url_template in tile_urls:
                tile_data = fetch_tile_with_retry(url_template, zoom, x, y, headers)
                if tile_data:
                    try:
                        img = Image.open(io.BytesIO(tile_data)).resize((tile_size * scale_factor, tile_size * scale_factor), Image.Resampling.LANCZOS)
                        stitched.paste(img, ((x - x_min) * tile_size * scale_factor, (y - y_min) * tile_size * scale_factor))
                        break
                    except Exception: pass
    
    def num2px(lat_deg, lon_deg, z):
        lat_rad = math.radians(lat_deg)
        n_tiles = 2.0 ** z
        return (lon_deg + 180.0) / 360.0 * n_tiles * tile_size * scale_factor, (1.0 - math.log(math.tan(lat_rad) + (1 / math.cos(lat_rad))) / math.pi) / 2.0 * n_tiles * tile_size * scale_factor
    
    px_w, py_n = num2px(n, w, zoom)
    px_e, py_s = num2px(s, e, zoom)
    base_x, base_y = x_min * tile_size * scale_factor, y_min * tile_size * scale_factor
    
    if not bounds_valid:
        pin_px_x, pin_px_y = num2px(pin_lat, pin_lon, zoom)
        left, top = int(pin_px_x - base_x - int(px_e - px_w) // 2), int(pin_px_y - base_y - int(py_s - py_n) // 2)
        right, bottom = left + int(px_e - px_w), top + int(py_s - py_n)
        cropped = stitched.crop((max(0, left), max(0, top), min(stitched.width, right), min(stitched.height, bottom))).convert("RGBA")
    else:
        cropped = stitched.crop((int(px_w - base_x), int(py_n - base_y), int(px_e - base_x), int(py_s - base_y))).convert("RGBA")
    
    draw = ImageDraw.Draw(cropped)
    pin_local_x = cropped.width // 2 if not bounds_valid else max(0, min(int(num2px(pin_lat, pin_lon, zoom)[0] - base_x) - int(px_w - base_x), cropped.width - 1))
    pin_local_y = cropped.height // 2 if not bounds_valid else max(0, min(int(num2px(pin_lat, pin_lon, zoom)[1] - base_y) - int(py_n - base_y), cropped.height - 1))
    
    radius = int((pin_size / 2) * scale_factor)
    draw.ellipse([pin_local_x - radius, pin_local_y - radius, pin_local_x + radius, pin_local_y + radius], fill=pin_color, outline=(255, 255, 255), width=2)
    
    star_size = int(radius * 0.55)
    star_points = []
    for i in range(10):
        angle = math.radians(i * 36 - 90)
        r = star_size if i % 2 == 0 else star_size * 0.4
        star_points.append((pin_local_x + r * math.cos(angle), pin_local_y + r * math.sin(angle)))
    draw.polygon(star_points, fill=(255, 255, 255))
    
    final_img = cropped.convert("RGB")
    img_byte_arr = io.BytesIO()
    final_img.save(img_byte_arr, format='PNG', quality=100)
    img_byte_arr.seek(0)
    return img_byte_arr

# --- ISOLATED MAP EDITOR ---
def render_isolated_map_editor():
    token_key = st.session_state.active_map_editor_token
    st.markdown('<div class="editor-card">', unsafe_allow_html=True)
    col_back, col_title = st.columns([1, 4])
    with col_back:
        if st.button("Back to Document", key="back_from_map"):
            st.session_state.restore_form_data = True
            st.session_state.active_map_editor_token = None
            st.rerun()
    with col_title: st.markdown(f"### Map Editor: {token_key}")
    st.markdown("</div><br>", unsafe_allow_html=True)

    style_key, coord_key, color_key, size_key, bounds_key, export_key = f"map_style_{token_key}", f"map_coord_{token_key}", f"map_color_{token_key}", f"map_size_{token_key}", f"map_bounds_{token_key}", f"map_exp_{token_key}"
    if style_key not in st.session_state: st.session_state[style_key] = "Satellite (Streets)"
    if coord_key not in st.session_state: st.session_state[coord_key] = "14.5995, 120.9842"
    if color_key not in st.session_state: st.session_state[color_key] = "#003366"
    if size_key not in st.session_state: st.session_state[size_key] = 20
    
    c_btn, c_style, c_color, c_size, c_coord = st.columns([1.4, 1.8, 0.8, 1.0, 2.8])
    with c_btn: export_clicked = st.button("Confirm and Export", type="primary", use_container_width=True)
    with c_style: basemap_style = st.selectbox("Basemap Layer", ["Satellite (Streets)", "Satellite (Labels)", "Satellite (Clean)", "Street Map", "OSM Carto Light", "Open Street Map"], key=style_key)
    with c_color: pin_color = st.color_picker("Pin", key=color_key)
    with c_size: pin_size = st.number_input("Size", min_value=8, max_value=64, key=size_key)
    with c_coord: coord_input = st.text_input("Coordinates", key=coord_key)
    
    try: plat, plon = map(float, coord_input.split(","))
    except ValueError: plat, plon = 14.5995, 120.9842

    if export_clicked:
        with st.spinner("Exporting Map..."):
            n, s, e, w = None, None, None, None
            if st.session_state.get(bounds_key):
                b = st.session_state[bounds_key]
                if b and "_northEast" in b and "_southWest" in b:
                    n, s, e, w = b["_northEast"]["lat"], b["_southWest"]["lat"], b["_northEast"]["lng"], b["_southWest"]["lng"]
            
            st.session_state[f"map_bytes_holder_{token_key}"] = generate_static_map_bounds(n, s, e, w, plat, plon, basemap_style, pin_color, int(pin_size))
            st.session_state.restore_form_data = True
            st.session_state.active_map_editor_token = None
            st.success("Map attached!")
            st.rerun()

    urls = get_tile_urls(basemap_style)
    m = folium.Map(location=[plat, plon], zoom_start=15, tiles=urls[0], attr=get_attribution(basemap_style))
    folium.Marker([plat, plon], draggable=True).add_to(m)
    Draw(export=False, draw_options={'polyline':False, 'polygon':False, 'circle':False, 'marker':False, 'circlemarker':False, 'rectangle':True}).add_to(m)
    
    map_data = st_folium(m, height=600, width=1300, use_container_width=True, key=f"int_map_{token_key}")
    if isinstance(map_data, dict):
        if map_data.get("bounds"): st.session_state[bounds_key] = map_data["bounds"]
        if map_data.get("last_marker_moved"):
            moved = map_data["last_marker_moved"]
            st.session_state[coord_key] = f"{round(moved['lat'], 5)}, {round(moved['lng'], 5)}"
            st.rerun()

# --- REPLACEMENT ENGINE ---
def smart_crop_to_fit(img_file, target_w_emu, target_h_emu):
    try:
        img = Image.open(img_file)
        target_ratio = target_w_emu / target_h_emu
        if img.size[0] / img.size[1] > target_ratio:
            new_w = int(img.size[1] * target_ratio)
            img = img.crop(((img.size[0] - new_w) // 2, 0, (img.size[0] - new_w) // 2 + new_w, img.size[1]))
        else:
            new_h = int(img.size[0] / target_ratio)
            img = img.crop((0, (img.size[1] - new_h) // 2, img.size[0], (img.size[1] - new_h) // 2 + new_h))
        buf = io.BytesIO()
        img.save(buf, format='PNG', quality=95)
        buf.seek(0)
        return buf
    except Exception: return img_file

def replace_text_in_paragraph(paragraph, text_inputs):
    """
    Scans and loops through target placeholders individually. 
    Maintains clean independent token evaluation without collapsing neighbor placeholders in the same text box.
    """
    for token, value in text_inputs.items():
        replacement = str(value) if value and str(value).strip() else ''
        
        # 1. Immediate full-run check to retain formatting
        for run in paragraph.runs:
            if token in run.text:
                run.text = run.text.replace(token, replacement)
        
        # 2. Aggressive multi-run fragmentation rebuild pass
        full_text = "".join(r.text for r in paragraph.runs)
        if token in full_text:
            updated_text = full_text.replace(token, replacement)
            if paragraph.runs:
                paragraph.runs[0].text = updated_text
                for secondary_run in paragraph.runs[1:]:
                    secondary_run.text = ""

def extract_placeholders_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    tokens, seen = [], set()
    for s in prs.slides:
        for shape in s.shapes:
            if shape.has_text_frame:
                for t in re.findall(r'\{\{.*?\}\}', shape.text):
                    if t not in seen: tokens.append(t); seen.add(t)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for t in re.findall(r'\{\{.*?\}\}', cell.text):
                            if t not in seen: tokens.append(t); seen.add(t)
    return tokens

def extract_placeholders_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    tokens, seen = [], set()
    for p in doc.paragraphs:
        for t in re.findall(r'\{\{.*?\}\}', p.text):
            if t not in seen: tokens.append(t); seen.add(t)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for t in re.findall(r'\{\{.*?\}\}', cell.text):
                    if t not in seen: tokens.append(t); seen.add(t)
    return tokens

def extract_placeholders(template_bytes, template_type):
    return extract_placeholders_from_pptx(template_bytes) if template_type == 'pptx' else extract_placeholders_from_docx(template_bytes)

def generate_pptx_bytes(template_bytes, text_inputs, image_inputs):
    prs = Presentation(io.BytesIO(template_bytes))
    for slide in prs.slides:
        shapes_to_delete, images_to_add = [], []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for img_token, img_file in image_inputs.items():
                    if img_token in shape.text and img_file is not None:
                        images_to_add.append((img_file, shape.left, shape.top, shape.width, shape.height))
                        shapes_to_delete.append(shape)
                        break
        for shape in slide.shapes:
            if shape not in shapes_to_delete:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs: replace_text_in_paragraph(paragraph, text_inputs)
                if hasattr(shape, 'table') and shape.table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs: replace_text_in_paragraph(paragraph, text_inputs)
        for img_file, left, top, width, height in images_to_add:
            try: slide.shapes.add_picture(smart_crop_to_fit(img_file, width, height), left, top, width=width, height=height)
            except Exception: pass
        for old_shape in shapes_to_delete:
            try:
                sp = old_shape._element
                sp.getparent().remove(sp)
            except Exception: pass
    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()

def generate_docx_bytes(template_bytes, text_inputs, image_inputs):
    doc = Document(io.BytesIO(template_bytes))
    for paragraph in doc.paragraphs:
        if not any(img_token in paragraph.text for img_token in image_inputs.keys()): replace_text_in_paragraph(paragraph, text_inputs)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.paragraphs:
                    for paragraph in cell.paragraphs: replace_text_in_paragraph(paragraph, text_inputs)
    stream = io.BytesIO()
    doc.save(stream)
    return stream.getvalue()

def get_download_filename(template_name, file_type):
    base_name = re.sub(r'^template_', '', template_name or "Document")
    base_name = re.sub(r'\.(pptx|docx)$', '', base_name)
    return f"Generated_{re.sub(r'[^\w\-_. ]', '_', base_name)}_{datetime.now().strftime('%m%d%Y')}.{file_type}"

def restore_form_data_from_session():
    if not st.session_state.saved_template_name: return False
    if any(f"val_{t}" in st.session_state for t in st.session_state.tokens): return True
    temp_path = get_temp_config_path(st.session_state.saved_template_name)
    if os.path.exists(temp_path):
        try:
            with open(temp_path, 'r', encoding='utf-8') as f:
                st.session_state.temp_form_data = json.load(f)
                for token, value in st.session_state.temp_form_data.items():
                    if st.session_state.custom_mapping.get(token, "Text") != "Image":
                        st.session_state[f"val_{token}"] = value
                return True
        except Exception: pass
    return False

def purge_all_temporary_data():
    if st.session_state.saved_template_name:
        temp_path = get_temp_config_path(st.session_state.saved_template_name)
        if os.path.exists(temp_path):
            try: os.remove(temp_path)
            except Exception: pass
    for token in st.session_state.tokens:
        st.session_state.pop(f"val_{token}", None)
        st.session_state.pop(f"map_bytes_holder_{token}", None)
    st.session_state.temp_form_data = {}

# --- ENTRY POINT HANDLER ---
if "custom_mapping" not in st.session_state: st.session_state.custom_mapping = {}
if "tokens" not in st.session_state: st.session_state.tokens = []
if "template_bytes" not in st.session_state: st.session_state.template_bytes = None
if "saved_template_name" not in st.session_state: st.session_state.saved_template_name = None
if "temp_form_data" not in st.session_state: st.session_state.temp_form_data = {}
if "restore_form_data" not in st.session_state: st.session_state.restore_form_data = False

if st.session_state.get("active_map_editor_token"):
    render_isolated_map_editor()
else:
    if st.session_state.restore_form_data:
        restore_form_data_from_session()
        st.session_state.restore_form_data = False

    st.markdown("<hr style='margin: 4px 0 12px 0;'>", unsafe_allow_html=True)
    st.markdown('<div class="workspace-card"><div class="section-header">Templates</div>', unsafe_allow_html=True)
    col_t1, col_t2 = st.columns(2)
    
    with col_t1:
        saved_templates = get_saved_templates()
        template_options = ["Select saved template"]
        for t in saved_templates: template_options.append(f"{t['display_name']} ({t['type']})")
        selected_template = st.selectbox("Load Template", template_options, key="saved_template_select", label_visibility="collapsed")
        
        if selected_template and selected_template != "Select saved template":
            t_display = selected_template.split(' (')[0].strip()
            for t in saved_templates:
                if t['display_name'] == t_display:
                    template_bytes = load_template_from_file(t['name'])
                    if template_bytes:
                        st.session_state.template_bytes = template_bytes
                        st.session_state.saved_template_name = t['name']
                        st.session_state.template_type = 'pptx' if t['name'].endswith('.pptx') else 'docx'
                        st.session_state.tokens = extract_placeholders(template_bytes, st.session_state.template_type)
                        config_data = load_config_from_file(t['name'].replace('.pptx', '').replace('.docx', '') + '_config.json')
                        if config_data: st.session_state.custom_mapping = config_data
                        restore_form_data_from_session()
                    break

    with col_t2:
        uploaded_template = st.file_uploader("Upload New Template", type=["pptx", "docx"], label_visibility="collapsed")
        if uploaded_template:
            st.session_state.template_bytes = uploaded_template.getvalue()
            st.session_state.saved_template_name = uploaded_template.name
            st.session_state.template_type = 'pptx' if uploaded_template.name.endswith('.pptx') else 'docx'
            st.session_state.tokens = extract_placeholders(st.session_state.template_bytes, st.session_state.template_type)
            st.session_state.temp_form_data = {}
            if st.button("Save Template", use_container_width=True):
                save_template_to_file(st.session_state.template_bytes, uploaded_template.name)
                st.rerun()

    if st.session_state.template_bytes is not None:
        st.markdown(f'<div class="saved-indicator">Active: {st.session_state.saved_template_name}</div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    text_data, image_data = {}, {}
    if st.session_state.template_bytes is not None and st.session_state.tokens:
        cta_sets = detect_cta_sets()
        if cta_sets:
            st.markdown("**Call to Action Presets**")
            cols = st.columns(len(cta_sets))
            for idx, cta_num in enumerate(sorted(cta_sets.keys())):
                with cols[idx]:
                    cta_name_token = cta_sets[cta_num]['tokens'].get('NAME')
                    curr_adv = st.session_state.get(f"val_{cta_name_token}", "")
                    sel_adv = st.selectbox(f"CTA {cta_num}", [""] + list(contacts_database.keys()), index=list(contacts_database.keys()).index(curr_adv)+1 if curr_adv in contacts_database else 0, key=f"cta_sel_{cta_num}")
                    if sel_adv and sel_adv != curr_adv:
                        apply_cta_preset_autofill(cta_num, sel_adv)
                        st.rerun()
        
        with st.expander("Data Type Mapping"):
            cols = st.columns(3)
            for idx, token in enumerate(st.session_state.tokens):
                with cols[idx % 3]:
                    curr_type = st.session_state.custom_mapping.get(token, "Text")
                    data_type = st.selectbox(f"Type for {token}", ["Text", "Image", "Map"], index=["Text", "Image", "Map"].index(curr_type))
                    if data_type != curr_type:
                        st.session_state.custom_mapping[token] = data_type
                        auto_save_config()
                        st.rerun()

        st.markdown('<div class="section-header">Placeholder Values</div>', unsafe_allow_html=True)
        for idx, token in enumerate(st.session_state.tokens):
            with (st.columns(2)[0] if idx % 2 == 0 else st.columns(2)[1]):
                curr_type = st.session_state.custom_mapping.get(token, "Text")
                st.markdown(f'<div class="placeholder-label">{token}</div>', unsafe_allow_html=True)
                
                if curr_type == "Image" and st.session_state.template_type == 'pptx':
                    image_data[token] = st.file_uploader(token, type=["png", "jpg"], key=f"val_{token}", label_visibility="collapsed")
                elif curr_type == "Map" and st.session_state.template_type == 'pptx':
                    saved_map = st.session_state.get(f"map_bytes_holder_{token}")
                    if saved_map:
                        image_data[token] = saved_map
                        st.caption("Map attached.")
                    
                    if st.button("Open Map Editor", key=f"btn_map_{token}", use_container_width=True):
                        st.session_state.active_map_editor_token = token
                        st.rerun()
                else:
                    curr_val = st.session_state.get(f"val_{token}", st.session_state.temp_form_data.get(token, ""))
                    new_val = st.text_input("", value=curr_val, key=f"val_{token}", label_visibility="collapsed")
                    if new_val != curr_val:
                        st.session_state.temp_form_data[token] = new_val
                        if st.session_state.saved_template_name:
                            with open(get_temp_config_path(st.session_state.saved_template_name), 'w') as f: json.dump(st.session_state.temp_form_data, f)
                    text_data[token] = new_val

        st.markdown('<div class="workspace-card"><div class="section-header">Download Document</div>', unsafe_allow_html=True)
        col1, col2 = st.columns(2)
        with col1:
            if st.session_state.template_type == 'pptx':
                st.download_button("Download PPTX", data=generate_pptx_bytes(st.session_state.template_bytes, text_data, image_data), file_name=get_download_filename(st.session_state.saved_template_name, "pptx"), use_container_width=True, on_click=purge_all_temporary_data)
        with col2:
            if st.session_state.template_type == 'docx':
                st.download_button("Download DOCX", data=generate_docx_bytes(st.session_state.template_bytes, text_data, image_data), file_name=get_download_filename(st.session_state.saved_template_name, "docx"), use_container_width=True, on_click=purge_all_temporary_data)
        st.markdown('</div>', unsafe_allow_html=True)
