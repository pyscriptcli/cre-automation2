import os
import io
import re
import json
import streamlit as st
from pptx import Presentation
from PIL import Image
from datetime import datetime
from docx import Document
import folium
from streamlit_folium import folium_static
import tempfile
import time
import base64
import requests
import math

# --- PROGRAMMATIC LIGHT MODE LOCK ---
_config_dir = ".streamlit"
_config_file = os.path.join(_config_dir, "config.toml")
if not os.path.exists(_config_file):
    os.makedirs(_config_dir, exist_ok=True)
    with open(_config_file, "w", encoding="utf-8") as f:
        f.write("[theme]\nbase=\"light\"\n")

# --- MINIMAL UI CSS ---
MINIMAL_CRE_SYSTEM = """
<style>
    #MainMenu {visibility: hidden;}
    header {visibility: hidden;}
    .stApp { margin-top: -50px; }
    .stDeployButton {display: none;}
    .stStatusWidget {display: none;}
    
    .stApp { background-color: #FFFFFF !important; color: #1A1A1A !important; font-family: 'Segoe UI', Arial, sans-serif !important; }
    div[data-testid="stHeader"] { background-color: #FFFFFF !important; display: none !important; }
    .block-container { padding-top: 0.5rem !important; padding-bottom: 0.5rem !important; max-width: 1200px !important; }
    
    div[data-baseweb="input"], div[data-baseweb="base-input"], div[role="textbox"], div[data-baseweb="select"], textarea {
        background-color: #FFFFFF !important; border: 1px solid #CCCCCC !important; border-radius: 4px !important;
        color: #1A1A1A !important;
    }
    div[data-baseweb="input"]:focus-within, div[data-baseweb="select"]:focus-within, textarea:focus { border-color: #003366 !important; box-shadow: none !important; }
    input[type="text"], .stTextInput input, div[data-baseweb="select"] div, textarea { color: #1A1A1A !important; font-size: 14px !important; }
    
    div[data-baseweb="select"] { min-height: 32px !important; }
    div[data-baseweb="select"] > div { min-height: 32px !important; padding: 0 8px !important; }
    div[data-baseweb="select"] select { font-size: 13px !important; padding: 2px 8px !important; }
    svg[data-testid="stSelectbox"] { width: 16px !important; height: 16px !important; }
    div[data-baseweb="select"] svg { width: 16px !important; height: 16px !important; }
    
    section[data-testid="stFileUploader"] { background-color: #F8F8F8 !important; border: 1px solid #CCCCCC !important; border-radius: 4px !important; padding: 4px 12px !important; }
    
    .workspace-card { background-color: #FFFFFF; border: 1px solid #E0E0E0; border-radius: 4px; padding: 16px; margin-bottom: 12px; }
    
    div.stButton > button { 
        background-color: #003366 !important; 
        color: #FFFFFF !important; 
        font-weight: 600 !important; 
        font-size: 11px !important; 
        border: none !important; 
        border-radius: 3px !important; 
        padding: 5px 12px !important; 
        width: 100% !important; 
        transition: background-color 0.15s ease; 
        min-height: 28px !important;
    }
    div.stButton > button:hover { 
        background-color: #002244 !important; 
        color: #FFFFFF !important; 
        transform: translateY(-1px);
        box-shadow: 0 2px 8px rgba(0, 51, 102, 0.3);
    }
    div.stButton > button:disabled { 
        background-color: #6688AA !important; 
        color: #CCCCCC !important; 
        cursor: not-allowed !important; 
    }
    
    div[data-testid="stDownloadButton"] > button { 
        background-color: #003366 !important;
        color: #FFFFFF !important;
        border-radius: 3px !important; 
        font-weight: 600 !important; 
        padding: 5px 12px !important; 
        width: 100% !important; 
        transition: all 0.15s ease;
        font-size: 11px !important;
        min-height: 28px !important;
    }
    div[data-testid="stDownloadButton"] > button:hover {
        background-color: #002244 !important;
        transform: translateY(-1px);
        box-shadow: 0 2px 8px rgba(0, 51, 102, 0.3);
    }
    
    .field-label { font-size: 13px !important; font-weight: 600 !important; color: #1A1A1A !important; padding-top: 6px; }
    .section-header { font-size: 15px !important; font-weight: 700 !important; color: #1A1A1A !important; margin-bottom: 10px; }
    .saved-indicator { background-color: #E8F5E9; padding: 6px 12px; border-radius: 4px; font-size: 13px; color: #2E7D32; border-left: 3px solid #2E7D32; margin-top: 6px; }
    
    .map-container { 
        border: 1px solid #E0E0E0; 
        border-radius: 4px; 
        padding: 8px; 
        background-color: #F8F9FA;
        margin: 8px 0;
    }
    .map-saved-indicator {
        background-color: #E3F2FD;
        padding: 4px 10px;
        border-radius: 4px;
        font-size: 11px;
        color: #003366;
        border-left: 3px solid #003366;
        margin: 4px 0;
    }
    .map-editor-header {
        background-color: #F8F9FA;
        padding: 8px 12px;
        border-radius: 4px;
        margin-bottom: 8px;
        border: 1px solid #E0E0E0;
        font-weight: 600;
        font-size: 13px;
    }
    .manual-capture-box {
        background-color: #FFF3E0;
        border: 2px dashed #FF9800;
        border-radius: 8px;
        padding: 16px;
        margin: 12px 0;
    }
    
    hr { margin: 12px 0 !important; border-color: #E0E0E0 !important; }
    .streamlit-expanderHeader { font-size: 14px !important; font-weight: 600 !important; }
</style>
"""

# --- FILE MANAGEMENT FUNCTIONS ---
def get_storage_dir():
    storage_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "stored_templates")
    os.makedirs(storage_dir, exist_ok=True)
    return storage_dir

def save_template_to_file(template_bytes, template_name):
    storage_dir = get_storage_dir()
    safe_name = re.sub(r'[^\w\-_. ]', '_', template_name)
    if not safe_name.endswith('.pptx') and not safe_name.endswith('.docx'):
        safe_name += '.docx'
    filepath = os.path.join(storage_dir, safe_name)
    with open(filepath, 'wb') as f:
        f.write(template_bytes)
    return filepath

def load_template_from_file(template_name):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        with open(filepath, 'rb') as f:
            return f.read()
    return None

def get_saved_templates():
    storage_dir = get_storage_dir()
    templates = []
    if os.path.exists(storage_dir):
        for file in os.listdir(storage_dir):
            if file.endswith('.pptx') or file.endswith('.docx'):
                filepath = os.path.join(storage_dir, file)
                stat = os.stat(filepath)
                templates.append({
                    'name': file,
                    'path': filepath,
                    'size': stat.st_size,
                    'modified': datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S'),
                    'type': 'PPTX' if file.endswith('.pptx') else 'DOCX'
                })
    return templates

def delete_template_file(template_name):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, template_name)
    if os.path.exists(filepath):
        os.remove(filepath)
        config_name = template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        config_path = os.path.join(storage_dir, config_name)
        if os.path.exists(config_path):
            os.remove(config_path)
        return True
    return False

def save_config_to_file(config_data, config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    serializable_config = {}
    for key, value in config_data.items():
        if isinstance(value, dict) and 'screenshot' in value:
            serializable_config[key] = {
                'type': value.get('type', 'Map'),
                'lat': value.get('lat'),
                'lng': value.get('lng'),
                'basemap': value.get('basemap', 'satellite')
            }
        else:
            serializable_config[key] = value
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(serializable_config, f, indent=4)
    return filepath

def load_config_from_file(config_name="template_config.json"):
    storage_dir = get_storage_dir()
    filepath = os.path.join(storage_dir, config_name)
    if os.path.exists(filepath):
        with open(filepath, 'r', encoding='utf-8') as f:
            return json.load(f)
    return None

def auto_save_config():
    if st.session_state.saved_template_name and st.session_state.custom_mapping:
        config_name = st.session_state.saved_template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
        save_config_to_file(st.session_state.custom_mapping, config_name)

# --- CORE UTILITIES ---
def smart_crop_to_fit(img_file, target_w_emu, target_h_emu):
    try:
        img = Image.open(img_file)
        img_w, img_h = img.size
        target_ratio = target_w_emu / target_h_emu
        img_ratio = img_w / img_h
        
        if img_ratio > target_ratio:
            new_w = int(img_h * target_ratio)
            left = (img_w - new_w) // 2
            img = img.crop((left, 0, left + new_w, img_h))
        else:
            new_h = int(img_w / target_ratio)
            top = (img_h - new_h) // 2
            img = img.crop((0, top, img_w, top + new_h))
            
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG')
        img_byte_arr.seek(0)
        return img_byte_arr
    except Exception:
        return img_file

def extract_placeholders_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    tokens = []
    seen = set()
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                found = re.findall(r'\{\{.*?\}\}', shape.text)
                for token in found:
                    if token not in seen:
                        tokens.append(token)
                        seen.add(token)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        found = re.findall(r'\{\{.*?\}\}', cell.text)
                        for token in found:
                            if token not in seen:
                                tokens.append(token)
                                seen.add(token)
    return tokens

def extract_placeholders_from_docx(docx_bytes):
    doc = Document(io.BytesIO(docx_bytes))
    tokens = []
    seen = set()
    
    for paragraph in doc.paragraphs:
        found = re.findall(r'\{\{.*?\}\}', paragraph.text)
        for token in found:
            if token not in seen:
                tokens.append(token)
                seen.add(token)
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                found = re.findall(r'\{\{.*?\}\}', cell.text)
                for token in found:
                    if token not in seen:
                        tokens.append(token)
                        seen.add(token)
    
    return tokens

def extract_placeholders(template_bytes, template_type):
    if template_type == 'pptx':
        return extract_placeholders_from_pptx(template_bytes)
    elif template_type == 'docx':
        return extract_placeholders_from_docx(template_bytes)
    return []

def replace_text_in_paragraph(paragraph, text_inputs):
    for run in paragraph.runs:
        for token, value in text_inputs.items():
            if token in run.text:
                replacement = str(value) if value else ''
                run.text = run.text.replace(token, replacement)
    
    if hasattr(paragraph, 'text') and paragraph.text:
        for token, value in text_inputs.items():
            if token in paragraph.text:
                if not paragraph.runs:
                    paragraph.add_run()
                for run in paragraph.runs:
                    if token in run.text:
                        replacement = str(value) if value else ''
                        run.text = run.text.replace(token, replacement)

def generate_pptx_bytes(template_bytes, text_inputs, image_inputs):
    prs = Presentation(io.BytesIO(template_bytes))
    
    for slide in prs.slides:
        shapes_to_delete = []
        images_to_add = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content = shape.text
                for img_token, img_file in image_inputs.items():
                    if img_token in text_content and img_file is not None:
                        images_to_add.append((img_file, shape.left, shape.top, shape.width, shape.height))
                        shapes_to_delete.append(shape)
                        break

        for shape in slide.shapes:
            if shape not in shapes_to_delete:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        replace_text_in_paragraph(paragraph, text_inputs)
                
                if hasattr(shape, 'table') and shape.table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    replace_text_in_paragraph(paragraph, text_inputs)

        for img_file, left, top, width, height in images_to_add:
            try:
                processed_img = smart_crop_to_fit(img_file, width, height)
                slide.shapes.add_picture(processed_img, left, top, width=width, height=height)
            except Exception:
                pass

        for old_shape in shapes_to_delete:
            try:
                sp = old_shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass

    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    return pptx_stream.getvalue()

def generate_docx_bytes(template_bytes, text_inputs, image_inputs):
    doc = Document(io.BytesIO(template_bytes))
    
    for paragraph in doc.paragraphs:
        has_image = False
        for img_token in image_inputs.keys():
            if img_token in paragraph.text:
                has_image = True
                break
        
        if not has_image:
            replace_text_in_paragraph(paragraph, text_inputs)
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    replace_text_in_paragraph(paragraph, text_inputs)
    
    doc_stream = io.BytesIO()
    doc.save(doc_stream)
    doc_stream.seek(0)
    return doc_stream.getvalue()

def get_download_filename(template_name, file_type):
    if template_name:
        base_name = re.sub(r'\.(pptx|docx)$', '', template_name)
        base_name = re.sub(r'[^\w\-_. ]', '_', base_name)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        return f"{base_name}_{timestamp}.{file_type}"
    else:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        return f"Generated_Document_{timestamp}.{file_type}"

# --- MAP FUNCTIONALITY - HYBRID AUTO + MANUAL ---
def get_basemap_tiles(basemap_choice):
    """Get the appropriate tile layer URL based on basemap choice"""
    basemaps = {
        'satellite': 'https://mt1.google.com/vt/lyrs=s&x={x}&y={y}&z={z}',
        'openstreetmap': 'https://{s}.tile.openstreetmap.org/{z}/{x}/{y}.png',
        'carto_light': 'https://{s}.basemaps.cartocdn.com/light_all/{z}/{x}/{y}{r}.png'
    }
    return basemaps.get(basemap_choice, basemaps['satellite'])

# --- AUTO CAPTURE METHODS ---

def capture_with_osm_tiles(lat, lng, basemap='satellite', zoom=15):
    """Method 1: Pure Python OSM tiles - No browser needed!"""
    try:
        import math
        from PIL import Image, ImageDraw
        
        def get_tile(zoom, x, y):
            url = f"https://tile.openstreetmap.org/{zoom}/{x}/{y}.png"
            response = requests.get(url, headers={"User-Agent": "OpenFlux/1.0"}, timeout=5)
            return Image.open(io.BytesIO(response.content))
        
        def lat_lon_to_tile(lat, lon, zoom):
            lat_rad = math.radians(lat)
            n = 2.0 ** zoom
            x = int((lon + 180.0) / 360.0 * n)
            y = int((1.0 - math.asinh(math.tan(lat_rad)) / math.pi) / 2.0 * n)
            return x, y
        
        # Only works with OSM basemap
        if basemap != 'openstreetmap':
            return None
        
        # Calculate tile coordinates
        center_x, center_y = lat_lon_to_tile(lat, lng, zoom)
        
        # Create canvas
        width, height = 800, 600
        tile_size = 256
        tiles_across = math.ceil(width / tile_size) + 1
        tiles_down = math.ceil(height / tile_size) + 1
        
        # Create combined image
        combined = Image.new('RGB', (tiles_across * tile_size, tiles_down * tile_size))
        
        # Download and stitch tiles
        for i in range(tiles_across):
            for j in range(tiles_down):
                tile_x = center_x - tiles_across//2 + i
                tile_y = center_y - tiles_down//2 + j
                try:
                    tile = get_tile(zoom, tile_x, tile_y)
                    combined.paste(tile, (i * tile_size, j * tile_size))
                except:
                    # Placeholder for missing tiles
                    placeholder = Image.new('RGB', (tile_size, tile_size), color='#E8ECF0')
                    combined.paste(placeholder, (i * tile_size, j * tile_size))
        
        # Crop to center
        crop_x = (combined.width - width) // 2
        crop_y = (combined.height - height) // 2
        cropped = combined.crop((crop_x, crop_y, crop_x + width, crop_y + height))
        
        # Draw pin at center
        draw = ImageDraw.Draw(cropped)
        pin_x, pin_y = width//2, height//2
        
        # Pin shadow
        draw.ellipse([pin_x-10, pin_y+20, pin_x+10, pin_y+30], fill='#B0B8C0')
        
        # Pin body
        draw.polygon([
            (pin_x, pin_y-20),
            (pin_x-12, pin_y+8),
            (pin_x+12, pin_y+8)
        ], fill='#FF0000')
        
        # Pin head
        draw.ellipse([pin_x-8, pin_y-8, pin_x+8, pin_y+8], fill='#FFFFFF')
        draw.ellipse([pin_x-4, pin_y-4, pin_x+4, pin_y+4], fill='#FF0000')
        
        # Save to bytes
        img_bytes = io.BytesIO()
        cropped.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes
        
    except Exception as e:
        print(f"OSM tiles capture failed: {e}")
        return None

def capture_with_google_static(lat, lng, basemap='satellite', zoom=15):
    """Method 2: Google Maps Static API"""
    try:
        api_key = "YOUR_GOOGLE_MAPS_API_KEY"
        if api_key == "YOUR_GOOGLE_MAPS_API_KEY":
            # Demo key - limited usage
            api_key = "AIzaSyA5oEohxJ-jB5WBR6pR3D8VtaY8X2CkT-8"
        
        maptype = 'satellite' if basemap == 'satellite' else 'roadmap'
        url = f"https://maps.googleapis.com/maps/api/staticmap?center={lat},{lng}&zoom={zoom}&size=800x600&maptype={maptype}&markers=color:red%7C{lat},{lng}&key={api_key}"
        
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            img = Image.open(io.BytesIO(response.content))
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            return img_bytes
        return None
    except:
        return None

def capture_with_osm_static(lat, lng, basemap='satellite', zoom=15):
    """Method 3: OSM Static API"""
    try:
        url = f"https://staticmap.openstreetmap.de/staticmap.php?center={lat},{lng}&zoom={zoom}&size=800x600&maptype=mapnik&markers={lat},{lng},red-pin"
        response = requests.get(url, timeout=10)
        if response.status_code == 200:
            img = Image.open(io.BytesIO(response.content))
            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            return img_bytes
        return None
    except:
        return None

def capture_with_selenium(lat, lng, basemap='satellite', zoom=15):
    """Method 4: Selenium with Chrome"""
    try:
        from selenium import webdriver
        from selenium.webdriver.chrome.options import Options
        from selenium.webdriver.chrome.service import Service
        from webdriver_manager.chrome import ChromeDriverManager
        
        tile_url = get_basemap_tiles(basemap)
        
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <link rel="stylesheet" href="https://unpkg.com/leaflet@1.9.4/dist/leaflet.css" />
            <script src="https://unpkg.com/leaflet@1.9.4/dist/leaflet.js"></script>
            <style>body, html {{ margin:0; padding:0; height:100%; }} #map {{ height:100vh; }}</style>
        </head>
        <body>
            <div id="map"></div>
            <script>
                var map = L.map('map').setView([{lat}, {lng}], {zoom});
                L.tileLayer('{tile_url}', {{maxZoom:20}}).addTo(map);
                var pinIcon = L.divIcon({{
                    html: `<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" width="32" height="32">
                        <path d="M12 2C8.13 2 5 5.13 5 9c0 5.25 7 13 7 13s7-7.75 7-13c0-3.87-3.13-7-7-7z" 
                              fill="#FF0000" stroke="#FFFFFF" stroke-width="1.5"/>
                        <circle cx="12" cy="9" r="2" fill="#FFFFFF"/>
                    </svg>`,
                    className:'', iconSize:[32,32], iconAnchor:[16,32]
                }});
                L.marker([{lat}, {lng}], {{icon: pinIcon, draggable:true}}).addTo(map);
            </script>
        </body>
        </html>
        """
        
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument('--disable-gpu')
        chrome_options.add_argument('--window-size=800,600')
        
        service = Service(ChromeDriverManager().install())
        driver = webdriver.Chrome(service=service, options=chrome_options)
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as f:
            f.write(html)
            html_path = f.name
        
        driver.get(f'file://{html_path}')
        time.sleep(2)
        screenshot = driver.get_screenshot_as_png()
        driver.quit()
        os.unlink(html_path)
        
        return io.BytesIO(screenshot)
        
    except Exception as e:
        print(f"Selenium capture failed: {e}")
        return None

def capture_with_playwright(lat, lng, basemap='satellite', zoom=15):
    """Method 5: Playwright"""
    try:
        from playwright.sync_api import sync_playwright
        
        tile_url = get_basemap_tiles(basemap)
        
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <link rel="stylesheet" href="https://unpkg.com/leaflet@1.9.4/dist/leaflet.css" />
            <script src="https://unpkg.com/leaflet@1.9.4/dist/leaflet.js"></script>
            <style>body, html {{ margin:0; padding:0; height:100%; }} #map {{ height:100vh; }}</style>
        </head>
        <body>
            <div id="map"></div>
            <script>
                var map = L.map('map').setView([{lat}, {lng}], {zoom});
                L.tileLayer('{tile_url}', {{maxZoom:20}}).addTo(map);
                var pinIcon = L.divIcon({{
                    html: `<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" width="32" height="32">
                        <path d="M12 2C8.13 2 5 5.13 5 9c0 5.25 7 13 7 13s7-7.75 7-13c0-3.87-3.13-7-7-7z" 
                              fill="#FF0000" stroke="#FFFFFF" stroke-width="1.5"/>
                        <circle cx="12" cy="9" r="2" fill="#FFFFFF"/>
                    </svg>`,
                    className:'', iconSize:[32,32], iconAnchor:[16,32]
                }});
                L.marker([{lat}, {lng}], {{icon: pinIcon, draggable:true}}).addTo(map);
            </script>
        </body>
        </html>
        """
        
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False) as f:
            f.write(html)
            html_path = f.name
        
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={'width': 800, 'height': 600})
            page.goto(f'file://{html_path}')
            page.wait_for_timeout(2000)
            screenshot = page.screenshot(full_page=True)
            browser.close()
        
        os.unlink(html_path)
        return io.BytesIO(screenshot)
        
    except Exception as e:
        print(f"Playwright capture failed: {e}")
        return None

def create_placeholder_image(lat, lng):
    """Ultimate fallback: Create a placeholder image"""
    try:
        from PIL import Image, ImageDraw
        
        img = Image.new('RGB', (800, 600), color='#F0F4F8')
        draw = ImageDraw.Draw(img)
        
        # Border
        draw.rectangle([10, 10, 790, 590], outline='#003366', width=2)
        
        # Grid
        for i in range(50, 800, 50):
            draw.line([(i, 10), (i, 590)], fill='#E0E5EC', width=1)
        for i in range(50, 600, 50):
            draw.line([(10, i), (790, i)], fill='#E0E5EC', width=1)
        
        # Pin
        pin_x, pin_y = 400, 250
        draw.ellipse([pin_x-10, pin_y+20, pin_x+10, pin_y+30], fill='#B0B8C0')
        draw.polygon([
            (pin_x, pin_y-20),
            (pin_x-12, pin_y+8),
            (pin_x+12, pin_y+8)
        ], fill='#FF0000')
        draw.ellipse([pin_x-8, pin_y-8, pin_x+8, pin_y+8], fill='#FFFFFF')
        draw.ellipse([pin_x-4, pin_y-4, pin_x+4, pin_y+4], fill='#FF0000')
        
        # Text
        draw.text((300, 400), f"Lat: {lat:.6f}, Lng: {lng:.6f}", fill='#003366')
        draw.text((350, 430), "Location Pin", fill='#003366')
        
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes
        
    except:
        # Super simple fallback
        img = Image.new('RGB', (800, 600), color='#FFFFFF')
        draw = ImageDraw.Draw(img)
        draw.text((300, 280), f"Location: {lat:.6f}, {lng:.6f}", fill='#000000')
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes

def capture_map_auto(lat, lng, basemap='satellite', zoom=15):
    """
    Try all automatic capture methods in order
    Returns: (image_bytes, method_used)
    """
    methods = [
        ('OSM Tiles (Pure Python)', capture_with_osm_tiles),
        ('Google Static API', capture_with_google_static),
        ('OSM Static API', capture_with_osm_static),
        ('Selenium', capture_with_selenium),
        ('Playwright', capture_with_playwright),
    ]
    
    for method_name, method in methods:
        try:
            result = method(lat, lng, basemap, zoom)
            if result is not None:
                return result, method_name
        except Exception as e:
            print(f"Method {method_name} failed: {e}")
            continue
    
    # Ultimate fallback
    return create_placeholder_image(lat, lng), "Placeholder Image"

def create_map_html_for_download(lat, lng, basemap='satellite', zoom=15):
    """Create HTML file for manual download"""
    tile_url = get_basemap_tiles(basemap)
    
    html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="utf-8">
        <title>Map Capture</title>
        <link rel="stylesheet" href="https://unpkg.com/leaflet@1.9.4/dist/leaflet.css" />
        <script src="https://unpkg.com/leaflet@1.9.4/dist/leaflet.js"></script>
        <style>
            body, html {{ margin: 0; padding: 0; height: 100%; font-family: Arial, sans-serif; }}
            #map {{ height: calc(100% - 50px); width: 100%; }}
            #controls {{
                height: 50px;
                background: #003366;
                color: white;
                display: flex;
                align-items: center;
                justify-content: center;
                gap: 20px;
                font-size: 14px;
            }}
            #controls span {{ opacity: 0.8; }}
            #controls strong {{ color: #FFD700; }}
            .info {{
                position: absolute;
                bottom: 70px;
                right: 20px;
                background: white;
                padding: 12px 16px;
                border-radius: 6px;
                box-shadow: 0 2px 10px rgba(0,0,0,0.2);
                font-family: monospace;
                font-size: 12px;
                z-index: 1000;
            }}
        </style>
    </head>
    <body>
        <div id="controls">
            <span>📍 Location: <strong>{lat:.6f}, {lng:.6f}</strong></span>
            <span>|</span>
            <span>🗺️ Basemap: <strong>{basemap}</strong></span>
            <span>|</span>
            <span>🔍 Zoom: <strong>{zoom}</strong></span>
        </div>
        <div id="map"></div>
        <div class="info">
            <b>📍 Pin Location</b><br>
            Lat: {lat:.6f}<br>
            Lng: {lng:.6f}<br>
            <span style="font-size:10px; color:#999;">Drag pin to adjust</span>
        </div>
        <script>
            var map = L.map('map').setView([{lat}, {lng}], {zoom});
            
            L.tileLayer('{tile_url}', {{
                maxZoom: 20,
                attribution: 'Map'
            }}).addTo(map);
            
            var pinIcon = L.divIcon({{
                html: `
                    <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" width="32" height="32">
                        <path d="M12 2C8.13 2 5 5.13 5 9c0 5.25 7 13 7 13s7-7.75 7-13c0-3.87-3.13-7-7-7z" 
                              fill="#FF0000" stroke="#FFFFFF" stroke-width="1.5"/>
                        <circle cx="12" cy="9" r="2" fill="#FFFFFF"/>
                    </svg>
                `,
                className: '',
                iconSize: [32, 32],
                iconAnchor: [16, 32]
            }});
            
            var marker = L.marker([{lat}, {lng}], {{
                icon: pinIcon,
                draggable: true
            }}).addTo(map);
            
            // Update info when dragged
            marker.on('dragend', function(e) {{
                var pos = marker.getLatLng();
                document.querySelector('.info').innerHTML = 
                    '<b>📍 Pin Location</b><br>Lat: ' + pos.lat.toFixed(6) + '<br>Lng: ' + pos.lng.toFixed(6) +
                    '<br><span style="font-size:10px; color:#999;">Drag pin to adjust</span>';
            }});
            
            // Instructions
            console.log('📸 To capture this map:');
            console.log('1. Adjust the pin and zoom as needed');
            console.log('2. Take a screenshot of the entire page');
            console.log('3. Upload the screenshot back to OpenFlux');
        </script>
    </body>
    </html>
    """
    return html

def map_editor_component(token, clean_label, default_lat=14.5995, default_lng=120.9842):
    """Map Editor with Hybrid Auto + Manual approach"""
    
    map_key = f"map_{token}"
    
    # Initialize session state for this map
    if map_key not in st.session_state:
        st.session_state[map_key] = {
            "lat": default_lat,
            "lng": default_lng,
            "screenshot": None,
            "saved": False,
            "basemap": "satellite",
            "zoom": 15,
            "editor_open": False,
            "auto_capture_failed": False,
            "capture_method": None
        }
    
    # Show current status
    if st.session_state[map_key]["saved"]:
        st.markdown(
            f'<div class="map-saved-indicator">✅ Location saved: {st.session_state[map_key]["lat"]:.6f}, {st.session_state[map_key]["lng"]:.6f}</div>', 
            unsafe_allow_html=True
        )
        if st.session_state[map_key]["screenshot"] is not None:
            st.image(st.session_state[map_key]["screenshot"], caption="Current Map", width=250)
            if st.session_state[map_key]["capture_method"]:
                st.caption(f"Captured via: {st.session_state[map_key]['capture_method']}")
    
    # Buttons
    col_btn, col_clear = st.columns([3, 1])
    with col_btn:
        if st.button("🗺️ Open Map Editor", key=f"open_editor_{token}", use_container_width=True):
            st.session_state[map_key]["editor_open"] = not st.session_state[map_key]["editor_open"]
            st.rerun()
    
    with col_clear:
        if st.button("Clear Map", key=f"clear_map_{token}", use_container_width=True):
            st.session_state[map_key]["saved"] = False
            st.session_state[map_key]["screenshot"] = None
            st.session_state[map_key]["auto_capture_failed"] = False
            st.rerun()
    
    # Map Editor Expander
    if st.session_state[map_key]["editor_open"]:
        with st.expander("🗺️ Map Editor", expanded=True):
            st.markdown('<div class="map-editor-header">📍 Set location, then click "Auto Capture" to try automatic capture</div>', unsafe_allow_html=True)
            
            # Editor controls
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                basemap_choice = st.selectbox(
                    "Basemap",
                    ["satellite", "openstreetmap", "carto_light"],
                    index=["satellite", "openstreetmap", "carto_light"].index(
                        st.session_state[map_key].get("basemap", "satellite")
                    ),
                    key=f"basemap_editor_{token}",
                    label_visibility="collapsed"
                )
                if basemap_choice != st.session_state[map_key].get("basemap", "satellite"):
                    st.session_state[map_key]["basemap"] = basemap_choice
            
            with col2:
                zoom = st.slider(
                    "Zoom",
                    min_value=10,
                    max_value=18,
                    value=st.session_state[map_key].get("zoom", 15),
                    key=f"zoom_editor_{token}"
                )
                st.session_state[map_key]["zoom"] = zoom
            
            with col3:
                st.markdown('<div style="padding-top: 24px;"></div>', unsafe_allow_html=True)
                if st.button("🔄 Refresh", key=f"refresh_map_{token}", use_container_width=True):
                    st.rerun()
            
            # Coordinate input
            current_lat = st.session_state[map_key]["lat"]
            current_lng = st.session_state[map_key]["lng"]
            
            default_coords = f"{current_lat:.6f}, {current_lng:.6f}"
            coords_input = st.text_input(
                "Coordinates (lat, lon)",
                value=default_coords,
                key=f"coords_{token}",
                help="Enter coordinates in format: lat, lon",
                placeholder="e.g., 14.5995, 120.9842"
            )
            
            # Parse coordinates
            lat, lng = parse_coordinates(coords_input)
            if lat is not None and lng is not None:
                if lat != st.session_state[map_key]["lat"] or lng != st.session_state[map_key]["lng"]:
                    st.session_state[map_key]["lat"] = lat
                    st.session_state[map_key]["lng"] = lng
                    st.rerun()
            
            # Display interactive map
            st.markdown('<div class="map-container">', unsafe_allow_html=True)
            
            try:
                current_lat = st.session_state[map_key]["lat"]
                current_lng = st.session_state[map_key]["lng"]
                
                tile_url = get_basemap_tiles(basemap_choice)
                
                m = folium.Map(
                    location=[current_lat, current_lng],
                    zoom_start=st.session_state[map_key]["zoom"],
                    width='100%',
                    height=400,
                    tiles=tile_url,
                    attr='Map'
                )
                
                pin_svg = """
                <svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 24 24" width="32" height="32">
                    <path d="M12 2C8.13 2 5 5.13 5 9c0 5.25 7 13 7 13s7-7.75 7-13c0-3.87-3.13-7-7-7z" 
                          fill="#FF0000" stroke="#FFFFFF" stroke-width="1.5"/>
                    <circle cx="12" cy="9" r="2" fill="#FFFFFF"/>
                </svg>
                """
                
                marker = folium.Marker(
                    [current_lat, current_lng],
                    popup=f"{clean_label}<br>{current_lat:.6f}, {current_lng:.6f}",
                    icon=folium.DivIcon(
                        html=pin_svg,
                        icon_size=(32, 32),
                        icon_anchor=(16, 32),
                        popup_anchor=(0, -32)
                    ),
                    draggable=True
                ).add_to(m)
                
                folium_static(m, width=700, height=400)
                st.caption("Drag the red pin to set location")
                
            except Exception as e:
                st.warning(f"Map display limited: {str(e)}")
                st.info("Enter coordinates manually and click Auto Capture")
            
            st.markdown('</div>', unsafe_allow_html=True)
            
            # --- AUTO CAPTURE SECTION ---
            st.markdown("---")
            col_auto, col_manual = st.columns([1, 1])
            
            with col_auto:
                if st.button("📷 Auto Capture", key=f"auto_capture_{token}", use_container_width=True):
                    with st.spinner("Trying automatic capture methods..."):
                        lat = st.session_state[map_key]["lat"]
                        lng = st.session_state[map_key]["lng"]
                        basemap = st.session_state[map_key]["basemap"]
                        zoom = st.session_state[map_key]["zoom"]
                        
                        result, method = capture_map_auto(lat, lng, basemap, zoom)
                        
                        if result is not None:
                            st.session_state[map_key]["screenshot"] = result
                            st.session_state[map_key]["saved"] = True
                            st.session_state[map_key]["capture_method"] = method
                            st.session_state[map_key]["auto_capture_failed"] = False
                            auto_save_config()
                            st.success(f"✅ Map captured successfully via: {method}")
                            st.rerun()
                        else:
                            st.session_state[map_key]["auto_capture_failed"] = True
                            st.error("❌ Auto capture failed. Please use manual capture below.")
                            st.rerun()
            
            with col_manual:
                if st.button("📄 Manual Capture", key=f"manual_capture_{token}", use_container_width=True):
                    # Generate HTML for download
                    lat = st.session_state[map_key]["lat"]
                    lng = st.session_state[map_key]["lng"]
                    basemap = st.session_state[map_key]["basemap"]
                    zoom = st.session_state[map_key]["zoom"]
                    
                    html_content = create_map_html_for_download(lat, lng, basemap, zoom)
                    st.session_state[f"manual_html_{token}"] = html_content
                    st.rerun()
            
            # --- MANUAL CAPTURE SECTION (shows after auto fails or manual clicked) ---
            if st.session_state[map_key]["auto_capture_failed"] or f"manual_html_{token}" in st.session_state:
                st.markdown('<div class="manual-capture-box">', unsafe_allow_html=True)
                st.markdown("### 📄 Manual Capture Instructions")
                st.markdown("""
                1. **Download the HTML file** below
                2. **Open it in your browser** (Chrome, Firefox, Safari)
                3. **Adjust the map** (drag pin, zoom) until it looks right
                4. **Take a screenshot** of the entire page
                5. **Upload the screenshot** using the uploader below
                """)
                
                # Download HTML button
                if f"manual_html_{token}" in st.session_state:
                    html_content = st.session_state[f"manual_html_{token}"]
                    st.download_button(
                        label="📥 Download Map HTML File",
                        data=html_content,
                        file_name=f"map_{token}.html",
                        mime="text/html",
                        use_container_width=True,
                        key=f"download_html_{token}"
                    )
                
                # Upload screenshot
                st.markdown("---")
                st.markdown("### 📤 Upload Your Screenshot")
                uploaded_file = st.file_uploader(
                    "Upload map screenshot (PNG or JPG)",
                    type=["png", "jpg", "jpeg"],
                    key=f"manual_upload_{token}",
                    label_visibility="collapsed"
                )
                
                if uploaded_file is not None:
                    # Process uploaded image
                    img = Image.open(uploaded_file)
                    img_bytes = io.BytesIO()
                    img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    
                    st.session_state[map_key]["screenshot"] = img_bytes
                    st.session_state[map_key]["saved"] = True
                    st.session_state[map_key]["capture_method"] = "Manual Upload"
                    st.session_state[map_key]["auto_capture_failed"] = False
                    auto_save_config()
                    
                    st.success("✅ Manual map uploaded successfully!")
                    st.rerun()
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            # Show captured preview if exists but not saved? (This is handled above)
    
    # Return screenshot if saved
    if st.session_state[map_key]["saved"] and st.session_state[map_key]["screenshot"] is not None:
        return st.session_state[map_key]["screenshot"]
    return None

def parse_coordinates(coord_string):
    """Parse coordinates from a string format: 'lat, lon'"""
    match = re.match(r'^\s*(-?\d+(?:\.\d+)?)\s*[,;]\s*(-?\d+(?:\.\d+)?)\s*$', coord_string.strip())
    if match:
        lat = float(match.group(1))
        lon = float(match.group(2))
        return lat, lon
    return None, None

def simple_uploader_row(label_text, allowed_types, key):
    st.markdown(f'<div class="field-label">{label_text}</div>', unsafe_allow_html=True)
    return st.file_uploader(label_text, type=allowed_types, key=f"val_{key}", label_visibility="collapsed")

# --- INIT APP ---
st.set_page_config(page_title="OpenFlux - Template Automation", layout="wide", initial_sidebar_state="collapsed")
st.markdown(MINIMAL_CRE_SYSTEM, unsafe_allow_html=True)

# Initialize session state
if "custom_mapping" not in st.session_state:
    st.session_state.custom_mapping = {}
if "tokens" not in st.session_state:
    st.session_state.tokens = []
if "template_bytes" not in st.session_state:
    st.session_state.template_bytes = None
if "saved_template_name" not in st.session_state:
    st.session_state.saved_template_name = None
if "template_loaded" not in st.session_state:
    st.session_state.template_loaded = False
if "template_type" not in st.session_state:
    st.session_state.template_type = None
if "delete_trigger" not in st.session_state:
    st.session_state.delete_trigger = False
if "show_delete_confirm" not in st.session_state:
    st.session_state.show_delete_confirm = False
if "template_to_delete" not in st.session_state:
    st.session_state.template_to_delete = None
if "save_success" not in st.session_state:
    st.session_state.save_success = False
if "saved_file_name" not in st.session_state:
    st.session_state.saved_file_name = None
if "clear_uploader" not in st.session_state:
    st.session_state.clear_uploader = False
if "map_data" not in st.session_state:
    st.session_state.map_data = {}

# --- MAIN UI ---
st.markdown("<hr style='margin: 4px 0 12px 0;'>", unsafe_allow_html=True)

# Template Management
st.markdown('<div class="workspace-card">', unsafe_allow_html=True)
st.markdown('<div class="section-header">Template Management</div>', unsafe_allow_html=True)

col_template1, col_template2 = st.columns(2)

with col_template1:
    saved_templates = get_saved_templates()
    template_options = ["Select saved template"]
    if saved_templates:
        for t in saved_templates:
            template_options.append(f"{t['name']} ({t['type']})")
    
    dropdown_col, delete_col = st.columns([4, 1])
    
    with dropdown_col:
        selected_template = st.selectbox(
            "Load Template",
            template_options,
            key="saved_template_select",
            label_visibility="collapsed"
        )
    
    with delete_col:
        if selected_template and selected_template != "Select saved template":
            template_name = selected_template.split(' (')[0]
            if st.button("Delete", key="delete_template", help="Delete this template"):
                st.session_state.show_delete_confirm = True
                st.session_state.template_to_delete = template_name
                st.rerun()
    
    if st.session_state.show_delete_confirm:
        st.warning(f"Are you sure you want to delete '{st.session_state.template_to_delete}'?")
        col_confirm1, col_confirm2 = st.columns([1, 1])
        with col_confirm1:
            if st.button("Yes, Delete", key="confirm_delete"):
                if delete_template_file(st.session_state.template_to_delete):
                    st.session_state.delete_trigger = True
                    st.session_state.template_bytes = None
                    st.session_state.saved_template_name = None
                    st.session_state.template_loaded = False
                    st.session_state.tokens = []
                    st.session_state.show_delete_confirm = False
                    st.session_state.template_to_delete = None
                    st.success(f"Deleted: {st.session_state.template_to_delete}")
                    st.rerun()
        with col_confirm2:
            if st.button("Cancel", key="cancel_delete"):
                st.session_state.show_delete_confirm = False
                st.session_state.template_to_delete = None
                st.rerun()
    
    if selected_template and selected_template != "Select saved template" and not st.session_state.delete_trigger:
        template_name = selected_template.split(' (')[0]
        template_bytes = load_template_from_file(template_name)
        if template_bytes:
            st.session_state.template_bytes = template_bytes
            st.session_state.saved_template_name = template_name
            st.session_state.template_loaded = True
            st.session_state.template_type = 'pptx' if template_name.endswith('.pptx') else 'docx'
            
            config_name = template_name.replace('.pptx', '').replace('.docx', '') + '_config.json'
            config_data = load_config_from_file(config_name)
            if config_data:
                st.session_state.custom_mapping = config_data
            
            tokens = extract_placeholders(template_bytes, st.session_state.template_type)
            st.session_state.tokens = tokens

with col_template2:
    uploader_key = "new_template_upload_clear" if st.session_state.clear_uploader else "new_template_upload"
    
    uploaded_template = st.file_uploader(
        "Upload New Template", 
        type=["pptx", "docx"], 
        label_visibility="collapsed", 
        key=uploader_key
    )
    
    if st.session_state.clear_uploader:
        st.session_state.clear_uploader = False
    
    if uploaded_template:
        template_bytes = uploaded_template.getvalue()
        st.session_state.template_bytes = template_bytes
        st.session_state.saved_template_name = None
        st.session_state.template_loaded = True
        st.session_state.template_type = 'pptx' if uploaded_template.name.endswith('.pptx') else 'docx'
        
        tokens = extract_placeholders(template_bytes, st.session_state.template_type)
        st.session_state.tokens = tokens
        
        if st.button("Save Template", key="save_template_btn", use_container_width=True):
            saved_path = save_template_to_file(template_bytes, uploaded_template.name)
            st.session_state.saved_template_name = uploaded_template.name
            
            if st.session_state.custom_mapping:
                config_name = uploaded_template.name.replace('.pptx', '').replace('.docx', '') + '_config.json'
                save_config_to_file(st.session_state.custom_mapping, config_name)
            
            st.session_state.save_success = True
            st.session_state.saved_file_name = uploaded_template.name
            st.session_state.clear_uploader = True
            st.rerun()

if st.session_state.save_success:
    st.success(f"Template '{st.session_state.saved_file_name}' saved successfully!")
    st.session_state.save_success = False
    st.session_state.saved_file_name = None

if st.session_state.template_bytes is not None:
    template_name = st.session_state.saved_template_name or "Unsaved Template"
    template_type = st.session_state.template_type or "Unknown"
    st.markdown(f'<div class="saved-indicator">Active: {template_name} ({template_type.upper()})</div>', unsafe_allow_html=True)

st.markdown('</div>', unsafe_allow_html=True)

# Placeholder Values
template_bytes = st.session_state.template_bytes
template_type = st.session_state.template_type
u_template = None
if template_bytes is not None:
    u_template = type('obj', (object,), {'getvalue': lambda: template_bytes})()

text_data = {}
image_data = {}
map_data = {}
field_types = {}

if u_template is not None and st.session_state.tokens:
    tokens = st.session_state.tokens
    
    if not tokens:
        st.info("No placeholders found in the template.")
    else:
        st.markdown('<div class="workspace-card">', unsafe_allow_html=True)
        st.markdown('<div class="section-header">Placeholder Values</div>', unsafe_allow_html=True)
        st.info(f"Found {len(tokens)} placeholders. Select type: Text, Image, or Map")
        
        mid_point = len(tokens) // 2
        col1, col2 = st.columns(2)
        
        with col1:
            for token in tokens[:mid_point]:
                clean_label = token.replace("{", "").replace("}", "")
                current_type = st.session_state.custom_mapping.get(token, "Text")
                col_a, col_b = st.columns([3, 1])
                
                with col_b:
                    st.markdown('<div style="padding-top: 6px;"></div>', unsafe_allow_html=True)
                    type_key = f"type_{token}"
                    data_type = st.selectbox(
                        "Type",
                        ["Text", "Image", "Map"],
                        index=0 if current_type == "Text" else (1 if current_type == "Image" else 2),
                        key=type_key,
                        label_visibility="collapsed"
                    )
                    if data_type != current_type:
                        st.session_state.custom_mapping[token] = data_type
                        auto_save_config()
                        st.rerun()
                
                with col_a:
                    if data_type == "Image" and template_type == 'pptx':
                        image_data[token] = simple_uploader_row(clean_label, ["png", "jpg", "jpeg"], token)
                        field_types[token] = "Image"
                        st.caption("Upload image (PNG, JPG)")
                    elif data_type == "Map":
                        st.session_state.map_data[token] = map_editor_component(token, clean_label)
                        field_types[token] = "Map"
                        st.caption("Click Open Map Editor to set location and capture map")
                    else:
                        if data_type == "Image" and template_type != 'pptx':
                            st.warning("Image replacement only supported in PPTX templates")
                        st.markdown(f'<div class="field-label">{clean_label}</div>', unsafe_allow_html=True)
                        text_data[token] = st.text_input(
                            clean_label, 
                            key=f"val_{token}", 
                            label_visibility="collapsed"
                        )
                        field_types[token] = "Text"
        
        with col2:
            for token in tokens[mid_point:]:
                clean_label = token.replace("{", "").replace("}", "")
                current_type = st.session_state.custom_mapping.get(token, "Text")
                col_a, col_b = st.columns([3, 1])
                
                with col_b:
                    st.markdown('<div style="padding-top: 6px;"></div>', unsafe_allow_html=True)
                    type_key = f"type_{token}_2"
                    data_type = st.selectbox(
                        "Type",
                        ["Text", "Image", "Map"],
                        index=0 if current_type == "Text" else (1 if current_type == "Image" else 2),
                        key=type_key,
                        label_visibility="collapsed"
                    )
                    if data_type != current_type:
                        st.session_state.custom_mapping[token] = data_type
                        auto_save_config()
                        st.rerun()
                
                with col_a:
                    if data_type == "Image" and template_type == 'pptx':
                        image_data[token] = simple_uploader_row(clean_label, ["png", "jpg", "jpeg"], token)
                        field_types[token] = "Image"
                        st.caption("Upload image (PNG, JPG)")
                    elif data_type == "Map":
                        st.session_state.map_data[token] = map_editor_component(token, clean_label)
                        field_types[token] = "Map"
                        st.caption("Click Open Map Editor to set location and capture map")
                    else:
                        if data_type == "Image" and template_type != 'pptx':
                            st.warning("Image replacement only supported in PPTX templates")
                        st.markdown(f'<div class="field-label">{clean_label}</div>', unsafe_allow_html=True)
                        text_data[token] = st.text_input(
                            clean_label, 
                            key=f"val_{token}", 
                            label_visibility="collapsed"
                        )
                        field_types[token] = "Text"
        
        st.markdown('</div>', unsafe_allow_html=True)

# Download Section
if u_template is not None:
    st.markdown('<div class="workspace-card">', unsafe_allow_html=True)
    st.markdown('<div class="section-header">Download Document</div>', unsafe_allow_html=True)
    
    # Merge map screenshots into image_data
    for token, map_screenshot in st.session_state.map_data.items():
        if map_screenshot is not None:
            image_data[token] = map_screenshot
    
    template_name = st.session_state.saved_template_name or "Generated_Document"
    base_template_name = re.sub(r'\.(pptx|docx)$', '', template_name)
    
    col1, col2 = st.columns(2)
    
    with col1:
        pptx_disabled = template_type != 'pptx'
        if pptx_disabled:
            st.button("Download PPTX", disabled=True, use_container_width=True, help="Only available for PPTX templates")
        else:
            try:
                pptx_data = generate_pptx_bytes(template_bytes, text_data, image_data)
                pptx_filename = get_download_filename(base_template_name, "pptx")
                st.download_button(
                    label="Download PPTX",
                    data=pptx_data,
                    file_name=pptx_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key="download_pptx"
                )
            except Exception as e:
                st.error(f"Error generating PPTX: {str(e)}")
    
    with col2:
        docx_disabled = template_type != 'docx'
        if docx_disabled:
            st.button("Download DOCX", disabled=True, use_container_width=True, help="Only available for DOCX templates")
        else:
            try:
                docx_data = generate_docx_bytes(template_bytes, text_data, image_data)
                if docx_data:
                    docx_filename = get_download_filename(base_template_name, "docx")
                    st.download_button(
                        label="Download DOCX",
                        data=docx_data,
                        file_name=docx_filename,
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        use_container_width=True,
                        key="download_docx"
                    )
                else:
                    st.error("Failed to generate document.")
            except Exception as e:
                st.error(f"Error generating document: {str(e)}")
    
    st.markdown('</div>', unsafe_allow_html=True)
else:
    st.info("Please upload or select a template to begin")

st.markdown("---")
st.caption("")
